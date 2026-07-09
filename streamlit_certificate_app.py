import os
import re
import zipfile
import shutil
import tempfile
import subprocess
import calendar
import hashlib
from datetime import date
from pathlib import Path
from typing import Tuple, Optional

import pandas as pd
import streamlit as st
from pptx import Presentation


# -----------------------------
# Page Configuration
# -----------------------------
st.set_page_config(
    page_title="Certificate Generator",
    page_icon="🎓",
    layout="wide",
    initial_sidebar_state="expanded",
)


# -----------------------------
# Professional UI Styling
# -----------------------------
def inject_custom_css() -> None:
    st.markdown(
        """
        <style>
            :root {
                --primary: #1E3A8A;
                --primary-soft: #DBEAFE;
                --accent: #B45309;
                --success: #15803D;
                --danger: #B91C1C;
                --text-dark: #111827;
                --text-muted: #6B7280;
                --card-bg: #FFFFFF;
                --page-bg: #F8FAFC;
                --border: #E5E7EB;
            }

            .stApp {
                background: linear-gradient(180deg, #F8FAFC 0%, #EEF2FF 100%);
            }

            section[data-testid="stSidebar"] {
                background: #FFFFFF;
                border-right: 1px solid #E5E7EB;
            }

            .block-container {
                padding-top: 1.8rem;
                padding-bottom: 3rem;
                max-width: 1220px;
            }

            .hero-card {
                background: linear-gradient(135deg, #0F172A 0%, #1E3A8A 55%, #92400E 100%);
                padding: 34px 36px;
                border-radius: 24px;
                color: white;
                box-shadow: 0 20px 45px rgba(15, 23, 42, 0.20);
                margin-bottom: 24px;
            }

            .hero-title {
                font-size: 38px;
                line-height: 1.15;
                font-weight: 800;
                margin: 0 0 10px 0;
                letter-spacing: -0.03em;
            }

            .hero-subtitle {
                font-size: 17px;
                color: #E5E7EB;
                max-width: 850px;
                margin-bottom: 0;
            }

            .section-card {
                background: var(--card-bg);
                border: 1px solid var(--border);
                border-radius: 20px;
                padding: 22px;
                box-shadow: 0 12px 30px rgba(15, 23, 42, 0.06);
                margin-bottom: 18px;
            }

            .mini-card {
                background: #FFFFFF;
                border: 1px solid #E5E7EB;
                border-radius: 18px;
                padding: 18px;
                min-height: 116px;
                box-shadow: 0 8px 18px rgba(15, 23, 42, 0.04);
            }

            .mini-card h4 {
                margin: 0 0 8px 0;
                color: #111827;
                font-size: 17px;
            }

            .mini-card p {
                margin: 0;
                color: #6B7280;
                font-size: 14px;
            }

            .step-badge {
                display: inline-flex;
                align-items: center;
                justify-content: center;
                width: 30px;
                height: 30px;
                border-radius: 999px;
                background: #DBEAFE;
                color: #1E3A8A;
                font-weight: 800;
                margin-right: 8px;
            }

            .status-pill {
                display: inline-block;
                padding: 6px 12px;
                border-radius: 999px;
                background: #ECFDF5;
                color: #166534;
                font-weight: 700;
                font-size: 13px;
                border: 1px solid #BBF7D0;
            }

            .muted-text {
                color: #6B7280;
                font-size: 14px;
            }

            div[data-testid="stMetric"] {
                background: #FFFFFF;
                border: 1px solid #E5E7EB;
                padding: 16px 18px;
                border-radius: 18px;
                box-shadow: 0 8px 18px rgba(15, 23, 42, 0.05);
            }

            div[data-testid="stMetricLabel"] p {
                color: #6B7280;
                font-size: 14px;
            }

            div[data-testid="stMetricValue"] {
                color: #111827;
                font-weight: 800;
            }

            .stButton > button {
                border-radius: 14px;
                padding: 0.75rem 1.1rem;
                font-weight: 800;
                border: 0;
            }

            .stDownloadButton > button {
                border-radius: 14px;
                font-weight: 700;
                border: 1px solid #D1D5DB;
                background: #FFFFFF;
            }

            .stDownloadButton > button:hover {
                border-color: #1E3A8A;
                color: #1E3A8A;
            }

            div[data-testid="stFileUploader"] {
                background: #FFFFFF;
                border: 1px dashed #CBD5E1;
                border-radius: 18px;
                padding: 14px;
            }

            div[data-testid="stDataFrame"] {
                border-radius: 16px;
                overflow: hidden;
            }

            hr {
                margin-top: 1.4rem;
                margin-bottom: 1.4rem;
            }

            .footer-note {
                background: #FFFBEB;
                border: 1px solid #FDE68A;
                color: #78350F;
                padding: 14px 16px;
                border-radius: 16px;
                font-size: 14px;
            }
        </style>
        """,
        unsafe_allow_html=True,
    )


# -----------------------------
# Helper functions
# -----------------------------
REQUIRED_COLUMNS = [
    "Name of Student",
    "Course",
    "Total Marks",
    "Grades",
    "Enrollment No.",
    "Exam Cycle",
]

PLACEHOLDERS = {
    "{{Name}}": "Name of Student",
    "{{SName}}": "Name of Student",
    "{{Course}}": "Course",
    "{{Marks}}": "Total Marks",
    "{{Grade}}": "Grades",
    "{{Enrollment}}": "Enrollment No.",
    "{{DateIssued}}": "Calculated from Exam Cycle",
}

MONTH_LOOKUP = {
    "JAN": 1, "JANUARY": 1,
    "FEB": 2, "FEBRUARY": 2,
    "MAR": 3, "MARCH": 3,
    "APR": 4, "APRIL": 4,
    "MAY": 5,
    "JUN": 6, "JUNE": 6,
    "JUL": 7, "JULY": 7,
    "AUG": 8, "AUGUST": 8,
    "SEP": 9, "SEPT": 9, "SEPTEMBER": 9,
    "OCT": 10, "OCTOBER": 10,
    "NOV": 11, "NOVEMBER": 11,
    "DEC": 12, "DECEMBER": 12,
}


def install_custom_fonts() -> None:
    """Install fonts from ./fonts on Linux/Streamlit Cloud. Windows uses installed system fonts."""
    fonts_source = Path("fonts")

    if not fonts_source.exists():
        return

    if os.name == "nt":
        return

    fonts_target = Path.home() / ".local" / "share" / "fonts"
    fonts_target.mkdir(parents=True, exist_ok=True)

    for pattern in ("*.ttf", "*.TTF", "*.otf", "*.OTF"):
        for font_file in fonts_source.glob(pattern):
            shutil.copy2(font_file, fonts_target / font_file.name)

    subprocess.run(
        ["fc-cache", "-f", "-v"],
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        check=False,
    )


def safe_filename(text: str) -> str:
    text = str(text).strip()
    text = re.sub(r'[\\/*?:"<>|]', "", text)
    text = re.sub(r"\s+", " ", text)
    return text or "certificate"


def format_certificate_date(date_value: date) -> str:
    return date_value.strftime("%d/%m/%Y")


def last_day_of_month(year: int, month: int) -> date:
    return date(year, month, calendar.monthrange(year, month)[1])


def get_last_completed_month_end(today: Optional[date] = None) -> date:
    today = today or date.today()
    if today.month == 1:
        return last_day_of_month(today.year - 1, 12)
    return last_day_of_month(today.year, today.month - 1)


def parse_exam_cycle(exam_cycle: str) -> tuple[int, int]:
    value = str(exam_cycle).strip()
    if not value:
        raise ValueError("Exam Cycle is blank")

    cleaned = re.sub(r"[^A-Za-z0-9]", "", value).upper()
    match = re.match(r"^([A-Z]+)(\d{2}|\d{4})$", cleaned)

    if not match:
        raise ValueError(f"Invalid Exam Cycle format: {exam_cycle}. Use values like Jan26, Apr26, Sep26.")

    month_text, year_text = match.groups()
    if month_text not in MONTH_LOOKUP:
        raise ValueError(f"Invalid month in Exam Cycle: {exam_cycle}")

    year = int(year_text)
    if len(year_text) == 2:
        year += 2000

    return MONTH_LOOKUP[month_text], year


def calculate_date_issued(exam_cycle: str, today: Optional[date] = None) -> str:
    today = today or date.today()
    month, year = parse_exam_cycle(exam_cycle)
    exam_cycle_end = last_day_of_month(year, month)
    last_completed_end = get_last_completed_month_end(today)
    issue_date = min(exam_cycle_end, last_completed_end)
    return format_certificate_date(issue_date)


def add_date_issued_preview(df: pd.DataFrame) -> pd.DataFrame:
    df = df.copy()
    df["Date Issued"] = df["Exam Cycle"].apply(calculate_date_issued)
    return df


def replace_text(shape, replacements: dict) -> None:
    """Replace all placeholders.
    Only {{Name}} uses Baguet Script.
    Other placeholders keep template font.
    """
    if not shape.has_text_frame:
        return

    FONT_NAME = "Baguet Script"

    for paragraph in shape.text_frame.paragraphs:
        full_text = "".join(run.text for run in paragraph.runs)

        if not full_text:
            continue

        updated = full_text
        name_found = "{{Name}}" in full_text

        for key, value in replacements.items():
            updated = updated.replace(key, str(value))

        if updated != full_text and paragraph.runs:
            first_run = paragraph.runs[0]
            first_run.text = updated

            if name_found:
                first_run.font.name = FONT_NAME

            for run in paragraph.runs[1:]:
                run.text = ""


def get_font_status() -> tuple[str, str]:
    """Return font status type and message without disturbing main UI."""
    if os.name == "nt":
        return (
            "info",
            "Windows: install Baguet Script by right-clicking the font file and choosing 'Install for all users'.",
        )

    try:
        result = subprocess.run(
            ["fc-match", "Baguet Script"],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            check=False,
        )
        matched = result.stdout.strip() or "No match returned"
        return "info", f"Linux font match: {matched}"
    except Exception as exc:
        return "warning", f"Could not check font status: {exc}"


def process_shape(shape, replacements: dict) -> None:
    if shape.shape_type == 6:  # Group shape
        for shp in shape.shapes:
            process_shape(shp, replacements)
        return

    replace_text(shape, replacements)


def create_zip(source_folder: Path, zip_path: Path) -> None:
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zipf:
        for file in sorted(source_folder.iterdir()):
            if file.is_file():
                zipf.write(file, arcname=file.name)


def convert_pptx_to_pdf(ppt_folder: Path, pdf_folder: Path) -> Tuple[bool, str]:
    pdf_folder.mkdir(parents=True, exist_ok=True)

    possible_paths = [
        shutil.which("soffice"),
        shutil.which("libreoffice"),
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]

    soffice = None
    for path in possible_paths:
        if path and os.path.exists(path):
            soffice = path
            break

    if soffice is None:
        return (
            False,
            "LibreOffice was not found. Install LibreOffice or add soffice/libreoffice to PATH.",
        )

    ppt_files = list(ppt_folder.glob("*.pptx"))
    if not ppt_files:
        return False, "No PPTX files found."

    errors = []
    converted = 0

    for ppt_file in ppt_files:
        command = [
            soffice,
            "--headless",
            "--nologo",
            "--nofirststartwizard",
            "--convert-to",
            "pdf",
            "--outdir",
            str(pdf_folder),
            str(ppt_file),
        ]

        result = subprocess.run(
            command,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
        )

        expected_pdf = pdf_folder / f"{ppt_file.stem}.pdf"
        if expected_pdf.exists():
            converted += 1
        else:
            errors.append(
                f"{ppt_file.name}\n"
                f"Return Code: {result.returncode}\n"
                f"STDOUT: {result.stdout}\n"
                f"STDERR: {result.stderr}"
            )

    if converted == 0:
        return False, "PDF conversion failed.\n\n" + "\n\n".join(errors[:3])

    if errors:
        return True, f"{converted} PDFs created. Some files failed:\n" + "\n\n".join(errors[:3])

    return True, f"Successfully created {converted} PDF(s)."


def generate_certificates(
    df: pd.DataFrame,
    template_path: Path,
    output_folder: Path,
    progress_bar=None,
    status_box=None,
) -> tuple[int, int, list[dict]]:
    output_folder.mkdir(parents=True, exist_ok=True)

    passed = 0
    failed = 0
    failed_students = []
    total = len(df)

    if total == 0:
        return 0, 0, []

    for count, (_, row) in enumerate(df.iterrows(), start=1):
        grade = str(row["Grades"]).strip().upper()

        if grade == "F":
            failed += 1
            failed_students.append(row.to_dict())
            if status_box:
                status_box.info(f"Skipped: {row['Name of Student']} because Grade is F")
        else:
            prs = Presentation(str(template_path))
            date_issued = row.get("Date Issued") or calculate_date_issued(row["Exam Cycle"])

            replacements = {
                "{{Name}}": row["Name of Student"],
                "{{Course}}": row["Course"],
                "{{Marks}}": row["Total Marks"],
                "{{Grade}}": row["Grades"],
                "{{DateIssued}}": date_issued,
                "{{Enrollment}}": row["Enrollment No."],
                "{{SName}}": row["Name of Student"],
            }

            for slide in prs.slides:
                for shape in slide.shapes:
                    process_shape(shape, replacements)

            filename = safe_filename(f"{row['Enrollment No.']}_{row['Name of Student']}.pptx")
            prs.save(output_folder / filename)
            passed += 1

            if status_box:
                status_box.success(f"Generated: {filename} | Date Issued: {date_issued}")

        if progress_bar:
            progress_bar.progress(count / total)

    return passed, failed, failed_students


def files_to_bytes_dict(folder: Path, pattern: str) -> dict[str, bytes]:
    file_dict = {}
    for file in sorted(folder.glob(pattern)):
        if file.is_file():
            file_dict[file.name] = file.read_bytes()
    return file_dict


def download_state_key(key_prefix: str, file_name: str) -> str:
    digest = hashlib.md5(file_name.encode("utf-8")).hexdigest()
    return f"downloaded_{key_prefix}_{digest}"


def mark_downloaded(state_key: str) -> None:
    st.session_state[state_key] = True


def render_file_name(file_name: str, downloaded: bool) -> None:
    if downloaded:
        st.markdown(
            f"<div style='color:#15803D; font-weight:800;'>✅ {file_name}</div>",
            unsafe_allow_html=True,
        )
    else:
        st.markdown(f"<div style='font-weight:600; color:#111827;'>{file_name}</div>", unsafe_allow_html=True)


def render_individual_downloads(title: str, files_dict: dict[str, bytes], mime: str, key_prefix: str) -> None:
    if not files_dict:
        return

    expander_key = f"expander_open_{key_prefix}"
    st.toggle(f"Show {title}", key=expander_key)

    if not st.session_state.get(expander_key, False):
        return

    with st.expander(title, expanded=True):
        search = st.text_input("Search by student name or enrollment number", key=f"search_{key_prefix}")
        filtered_files = {
            name: data
            for name, data in files_dict.items()
            if search.lower() in name.lower()
        }

        st.caption(f"Showing {len(filtered_files)} of {len(files_dict)} files")

        if not filtered_files:
            st.warning("No matching certificate found.")
            return

        for i, (file_name, file_data) in enumerate(filtered_files.items(), start=1):
            state_key = download_state_key(key_prefix, file_name)
            downloaded = st.session_state.get(state_key, False)

            c1, c2 = st.columns([4, 1])
            with c1:
                render_file_name(file_name, downloaded)
            with c2:
                st.download_button(
                    label="Downloaded" if downloaded else "Download",
                    data=file_data,
                    file_name=file_name,
                    mime=mime,
                    key=f"{key_prefix}_{i}_{file_name}",
                    on_click=mark_downloaded,
                    args=(state_key,),
                )


def init_session_state() -> None:
    defaults = {
        "generated": False,
        "ppt_zip_bytes": None,
        "pdf_zip_bytes": None,
        "skipped_excel_bytes": None,
        "ppt_files": {},
        "pdf_files": {},
        "summary": {},
        "pdf_message": "",
    }
    for key, value in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = value


def clear_downloaded_statuses() -> None:
    for key in list(st.session_state.keys()):
        if str(key).startswith("downloaded_"):
            del st.session_state[key]


def clear_previous_results() -> None:
    st.session_state.generated = False
    st.session_state.ppt_zip_bytes = None
    st.session_state.pdf_zip_bytes = None
    st.session_state.skipped_excel_bytes = None
    st.session_state.ppt_files = {}
    st.session_state.pdf_files = {}
    st.session_state.summary = {}
    st.session_state.pdf_message = ""
    clear_downloaded_statuses()


def render_sidebar(convert_pdf_default: bool = False) -> bool:
    with st.sidebar:
        st.markdown("### ⚙️ Generation Settings")
        convert_pdf = st.checkbox("Create PDF copies also", value=convert_pdf_default)

        st.markdown("---")
        st.markdown("### 🧩 Template Placeholders")
        st.caption("Use these exact placeholders inside your PPT template.")
        st.code(
            "{{Name}}\n{{Course}}\n{{Marks}}\n{{Grade}}\n{{DateIssued}}\n{{Enrollment}}\n{{SName}}",
            language="text",
        )

        st.markdown("---")
        st.markdown("### 📅 Date Rule")
        st.info(
            "Date Issued is calculated from `Exam Cycle`. Example: Jan26 → 31/01/2026. "
            "Current/future cycles use the last completed month."
        )

        st.markdown("---")
        st.markdown("### 🔤 Font Status")
        font_type, font_message = get_font_status()
        if font_type == "warning":
            st.warning(font_message)
        else:
            st.caption(font_message)

        st.markdown("---")
        st.caption("For Streamlit Cloud PDF conversion, keep `libreoffice` inside packages.txt.")

    return convert_pdf


def render_hero() -> None:
    st.markdown(
        """
        <div class="hero-card">
            <div class="status-pill">Certificate Automation</div>
            <h1 class="hero-title">Excel to PPT Certificate Generator</h1>
            <p class="hero-subtitle">
                Upload student data and your PowerPoint template. The app generates personalized PPT certificates,
                optional PDF copies, bulk ZIP files, skipped-student reports, and individual certificate downloads.
            </p>
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_workflow_cards() -> None:
    c1, c2, c3 = st.columns(3)
    with c1:
        st.markdown(
            """
            <div class="mini-card">
                <h4><span class="step-badge">1</span>Upload Excel</h4>
                <p>Student data should contain all required columns including Exam Cycle.</p>
            </div>
            """,
            unsafe_allow_html=True,
        )
    with c2:
        st.markdown(
            """
            <div class="mini-card">
                <h4><span class="step-badge">2</span>Upload Template</h4>
                <p>Your PPT template should contain the supported placeholders.</p>
            </div>
            """,
            unsafe_allow_html=True,
        )
    with c3:
        st.markdown(
            """
            <div class="mini-card">
                <h4><span class="step-badge">3</span>Generate & Download</h4>
                <p>Download certificates in bulk or individually after generation.</p>
            </div>
            """,
            unsafe_allow_html=True,
        )


# -----------------------------
# Streamlit App
# -----------------------------
inject_custom_css()
init_session_state()
install_custom_fonts()

convert_pdf = render_sidebar()
render_hero()
render_workflow_cards()

st.markdown("### Upload Files")
upload_col1, upload_col2 = st.columns(2)

with upload_col1:
    st.markdown("**Student Excel File**")
    excel_file = st.file_uploader(
        "Upload Excel File",
        type=["xlsx"],
        label_visibility="collapsed",
        help="Required columns: Name of Student, Course, Total Marks, Grades, Enrollment No., Exam Cycle",
    )

with upload_col2:
    st.markdown("**PowerPoint Certificate Template**")
    ppt_template = st.file_uploader(
        "Upload PPT Template",
        type=["pptx"],
        label_visibility="collapsed",
        help="The PPT template must contain placeholders like {{Name}}, {{Course}}, {{DateIssued}} etc.",
    )

st.markdown("---")

preview_df = None
if excel_file:
    try:
        preview_df = pd.read_excel(excel_file).fillna("")

        missing_columns = [col for col in REQUIRED_COLUMNS if col not in preview_df.columns]
        if missing_columns:
            st.error("Missing required columns: " + ", ".join(missing_columns))
            st.markdown("### Excel Preview")
            st.dataframe(preview_df.head(10), use_container_width=True)
        else:
            try:
                preview_df = add_date_issued_preview(preview_df)
            except Exception as e:
                st.error(f"Could not calculate Date Issued from Exam Cycle: {e}")

            total_students = len(preview_df)
            pass_count = len(preview_df[preview_df["Grades"].astype(str).str.strip().str.upper() != "F"])
            fail_count = total_students - pass_count

            st.markdown("### Data Summary")
            m1, m2, m3 = st.columns(3)
            m1.metric("Total Students", total_students)
            m2.metric("Certificates to Generate", pass_count)
            m3.metric("Skipped Grade F", fail_count)

            st.markdown("### Excel Preview")
            st.dataframe(preview_df.head(10), use_container_width=True)

    except Exception as e:
        st.error(f"Could not read Excel file: {e}")
        preview_df = None
else:
    st.info("Upload an Excel file to preview student records and validate required columns.")

button_disabled = not excel_file or not ppt_template or preview_df is None

button_col1, button_col2 = st.columns([1, 3])
with button_col1:
    generate_clicked = st.button(
        "🚀 Generate Certificates",
        disabled=button_disabled,
        type="primary",
        use_container_width=True,
    )
with button_col2:
    if button_disabled:
        st.caption("Upload both files and ensure all required Excel columns are available.")
    else:
        st.caption("Ready to generate certificates.")

if generate_clicked:
    clear_previous_results()

    with tempfile.TemporaryDirectory() as tmpdir:
        tmpdir = Path(tmpdir)
        excel_path = tmpdir / "students.xlsx"
        template_path = tmpdir / "template.pptx"
        output_folder = tmpdir / "Certificates"
        pdf_folder = tmpdir / "PDF_Certificates"
        ppt_zip = tmpdir / "Certificates.zip"
        pdf_zip = tmpdir / "PDF_Certificates.zip"
        skipped_file = tmpdir / "Skipped_Students.xlsx"

        excel_path.write_bytes(excel_file.getvalue())
        template_path.write_bytes(ppt_template.getvalue())

        df = pd.read_excel(excel_path).fillna("")
        missing_columns = [col for col in REQUIRED_COLUMNS if col not in df.columns]

        if missing_columns:
            st.error("Cannot generate. Missing columns: " + ", ".join(missing_columns))
        else:
            try:
                df = add_date_issued_preview(df)
            except Exception as e:
                st.error(f"Cannot generate. Invalid Exam Cycle data: {e}")
                st.stop()

            st.markdown("### Generation Progress")
            progress_bar = st.progress(0)
            status_box = st.empty()

            passed, failed, failed_students = generate_certificates(
                df=df,
                template_path=template_path,
                output_folder=output_folder,
                progress_bar=progress_bar,
                status_box=status_box,
            )

            st.session_state.ppt_files = files_to_bytes_dict(output_folder, "*.pptx")

            create_zip(output_folder, ppt_zip)
            st.session_state.ppt_zip_bytes = ppt_zip.read_bytes()

            if failed_students:
                pd.DataFrame(failed_students).to_excel(skipped_file, index=False)
                st.session_state.skipped_excel_bytes = skipped_file.read_bytes()

            if convert_pdf:
                with st.spinner("Converting PPT files to PDF..."):
                    success, message = convert_pptx_to_pdf(output_folder, pdf_folder)
                st.session_state.pdf_message = message

                if success and any(pdf_folder.glob("*.pdf")):
                    st.session_state.pdf_files = files_to_bytes_dict(pdf_folder, "*.pdf")
                    create_zip(pdf_folder, pdf_zip)
                    st.session_state.pdf_zip_bytes = pdf_zip.read_bytes()
                else:
                    st.session_state.pdf_files = {}
                    st.session_state.pdf_zip_bytes = None
            else:
                st.session_state.pdf_message = "PDF conversion was not selected."

            st.session_state.summary = {
                "total": len(df),
                "ppt_created": passed,
                "skipped": failed,
                "pdf_created": len(st.session_state.pdf_files),
            }
            st.session_state.generated = True
            st.rerun()


# -----------------------------
# Download Center
# -----------------------------
if st.session_state.generated:
    st.success("Certificate generation completed successfully.")

    st.markdown("### Final Summary")
    r1, r2, r3, r4 = st.columns(4)
    r1.metric("Total Students", st.session_state.summary.get("total", 0))
    r2.metric("PPT Certificates", st.session_state.summary.get("ppt_created", 0))
    r3.metric("PDF Certificates", st.session_state.summary.get("pdf_created", 0))
    r4.metric("Skipped", st.session_state.summary.get("skipped", 0))

    st.markdown("---")
    st.markdown("### Download Center")
    st.caption("Download all certificates together or search and download individual files.")

    st.markdown("#### Bulk Downloads")
    z1, z2, z3 = st.columns(3)

    with z1:
        if st.session_state.ppt_zip_bytes:
            st.download_button(
                label="⬇️ PPT Certificates ZIP",
                data=st.session_state.ppt_zip_bytes,
                file_name="Certificates.zip",
                mime="application/zip",
                key="download_ppt_zip",
                use_container_width=True,
            )

    with z2:
        if st.session_state.pdf_zip_bytes:
            st.download_button(
                label="⬇️ PDF Certificates ZIP",
                data=st.session_state.pdf_zip_bytes,
                file_name="PDF_Certificates.zip",
                mime="application/zip",
                key="download_pdf_zip",
                use_container_width=True,
            )
        elif convert_pdf:
            st.warning(st.session_state.pdf_message or "PDF ZIP not available.")
        else:
            st.info("PDF conversion not selected.")

    with z3:
        if st.session_state.skipped_excel_bytes:
            st.download_button(
                label="⬇️ Skipped Students Excel",
                data=st.session_state.skipped_excel_bytes,
                file_name="Skipped_Students.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                key="download_skipped_excel",
                use_container_width=True,
            )
        else:
            st.info("No skipped students report.")

    st.markdown("---")
    st.markdown("#### Individual Downloads")

    render_individual_downloads(
        title="Individual PPT Certificates",
        files_dict=st.session_state.ppt_files,
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        key_prefix="ppt_individual",
    )

    render_individual_downloads(
        title="Individual PDF Certificates",
        files_dict=st.session_state.pdf_files,
        mime="application/pdf",
        key_prefix="pdf_individual",
    )

    if not st.session_state.pdf_files:
        st.info("Individual PDF downloads will appear only when PDF conversion is selected and successful.")

st.markdown("---")
st.markdown(
    """
    <div class="footer-note">
        <b>Deployment Note:</b> For PDF conversion on Streamlit Cloud, add <code>libreoffice</code> in <code>packages.txt</code>. 
        For custom fonts, keep your .otf/.ttf font files inside a <code>fonts</code> folder in the GitHub repository.
    </div>
    """,
    unsafe_allow_html=True,
)
